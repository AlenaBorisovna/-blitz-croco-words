# Поменяйте код так, чтобы он выводил текст каждого шэйпа на экран.
# Сохраните из этой презентации только слова, всякие "Правила", "Тур 1" и другие объяснение не нужны.
# Сохраните все слова в файл words.txt
# Закомитьте его тоже.


from pptx import Presentation


def get_file():
    prs = Presentation("Osennyaya_igra_11.pptx")

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            print(shape.text)


if __name__ == '__main__':
    get_file()
