import os
from zipfile import ZipFile
from pathlib import Path
from typing import IO
from pptx import Presentation
from pyaspeller import YandexSpeller

FILE_TXT: str = 'words.txt'
FILE_ZIP = 'croco-blitz-source.zip'


def is_not_valid(text: str) -> bool:
    return ' ' in text or ':' in text or '-' in text


def get_words_presentation(file: IO[bytes]) -> set[str]:
    prs = Presentation(file)
    # words_list = list()
    words_set = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if is_not_valid(shape.text):
                continue
            words_set.add(shape.text)
    return words_set


def write_txt(filename: str, words: set[str] or str):
    try:
        with open(filename, 'w', encoding='utf-8') as file_text:
            for word in words:
                file_text.write(word + '\n')
        print(f"Слова записаны в файл: {filename}")
    except IOError as e:
        print(f"error: {e}")


def speller_word(words: set[str]) -> list[str]:
    print(f"проверка правильности написания слов {len(words)} с помощью Yandex Speller")
    spellers = YandexSpeller()
    fixes = spellers.spelled(' '.join(words))
    print('проверка орфографии выполнена')
    return fixes.split(' ')
